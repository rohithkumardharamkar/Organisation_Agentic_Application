import os
import hashlib
import math
from typing import Dict, Any, List, Optional
import pypdf
import docx
import pptx
from qdrant_client import QdrantClient
from qdrant_client.http import models as qmodels

# Initialize Qdrant client in persistent local mode
QDRANT_DB_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data", "qdrant_db")
os.makedirs(os.path.dirname(QDRANT_DB_PATH), exist_ok=True)

try:
    qdrant_client = QdrantClient(path=QDRANT_DB_PATH)
except Exception as e:
    print(f"Warning: Failed to lock persistent Qdrant database ({e}). Falling back to in-memory Qdrant client.")
    qdrant_client = QdrantClient(location=":memory:")

COLLECTION_NAME = "knowledge_documents"

# Ensure the collection is initialized with the correct configuration
try:
    qdrant_client.get_collection(COLLECTION_NAME)
except Exception:
    qdrant_client.create_collection(
        collection_name=COLLECTION_NAME,
        vectors_config=qmodels.VectorParams(size=1536, distance=qmodels.Distance.COSINE)
    )

def generate_semantic_bow_embedding(text: str) -> List[float]:
    """
    Generates a high-fidelity 1536-dimensional Bag-of-Words dense vector.
    Deterministically maps words to indices using MD5 hashes and normalizes the vector.
    This creates a dense representation where cosine similarity corresponds to term overlap.
    """
    vector = [0.0] * 1536
    # Tokenize and clean text
    words = [w.strip(".,!?;:()\"'[]{}") for w in text.lower().split() if w.strip()]
    
    # We weight words dynamically (tf-like weighting)
    for w in words:
        if len(w) < 2:
            continue
        # Project word hash into 1536 dimensions
        idx = int(hashlib.md5(w.encode("utf-8")).hexdigest(), 16) % 1536
        vector[idx] += 1.0
        
    # L2 Normalization
    sq_sum = sum(x * x for x in vector)
    norm = math.sqrt(sq_sum)
    if norm > 0:
        vector = [x / norm for x in vector]
    else:
        # Fallback dummy vector if text is empty
        vector[0] = 1.0
    return vector

class QdrantService:
    @staticmethod
    def parse_document(filepath: str) -> str:
        """Parses PDF, Word, PowerPoint, or text files into a raw string."""
        ext = os.path.splitext(filepath)[1].lower()
        text_content = []
        
        if ext == ".pdf":
            reader = pypdf.PdfReader(filepath)
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    text_content.append(t)
                    
        elif ext == ".docx":
            doc = docx.Document(filepath)
            for para in doc.paragraphs:
                if para.text:
                    text_content.append(para.text)
                    
        elif ext == ".pptx":
            prs = pptx.Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content.append(shape.text)
                        
        else: # Fallback to txt/raw parsing
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
                
        return "\n".join(text_content)

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 800, overlap: int = 150) -> List[str]:
        """Splits raw text into overlapping paragraph-friendly chunks."""
        chunks = []
        if not text:
            return chunks
            
        # Clean text
        text = " ".join(text.split())
        
        start = 0
        text_len = len(text)
        
        while start < text_len:
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk)
            start += chunk_size - overlap
            
        return chunks

    @classmethod
    async def index_document(
        cls,
        filename: str,
        filepath: str,
        category: str,
        uploaded_by: str,
        allowed_roles: List[str]
    ) -> Dict[str, Any]:
        """Parses, chunks, embeds, and indexes a file into Qdrant."""
        try:
            # 1. Parse text
            raw_text = cls.parse_document(filepath)
            if not raw_text.strip():
                raise ValueError("Parsed document contains no text.")
                
            # 2. Chunk text
            chunks = cls.chunk_text(raw_text)
            
            # 3. Create Points for Qdrant
            points = []
            for i, chunk in enumerate(chunks):
                vector = generate_semantic_bow_embedding(chunk)
                point_id = hashlib.md5(f"{filename}_{i}".encode("utf-8")).hexdigest()
                
                payload = {
                    "filename": filename,
                    "filepath": filepath,
                    "category": category,
                    "uploaded_by": uploaded_by,
                    "allowed_roles": allowed_roles,
                    "content": chunk,
                    "chunk_index": i,
                    "version": 1
                }
                
                points.append(
                    qmodels.PointStruct(
                        id=point_id,
                        vector=vector,
                        payload=payload
                    )
                )
                
            # 4. Upsert to Qdrant
            qdrant_client.upsert(
                collection_name=COLLECTION_NAME,
                points=points
            )
            
            return {
                "status": "success",
                "chunks_count": len(chunks),
                "message": f"Successfully indexed '{filename}' with {len(chunks)} chunks."
            }
        except Exception as e:
            return {
                "status": "error",
                "message": f"Indexing failed: {str(e)}"
            }

    @staticmethod
    async def search(
        query: str,
        user_role: str,
        category: Optional[str] = None,
        limit: int = 5
    ) -> List[Dict[str, Any]]:
        """Queries Qdrant vector store with role-based filtering and optional category filter."""
        query_vector = generate_semantic_bow_embedding(query)
        
        # Enforce Role-Based Access Control Filters
        filter_conditions = []
        
        if user_role != "HR":
            # Non-HR roles can only view documents containing their role or labeled 'all'
            role_condition = qmodels.Filter(
                should=[
                    qmodels.FieldCondition(
                        key="allowed_roles",
                        match=qmodels.MatchValue(value=user_role)
                    ),
                    qmodels.FieldCondition(
                        key="allowed_roles",
                        match=qmodels.MatchValue(value="all")
                    ),
                    # Support cases where allowed_roles includes 'Employee' as a baseline fallback
                    qmodels.FieldCondition(
                        key="allowed_roles",
                        match=qmodels.MatchValue(value="Employee")
                    )
                ]
            )
            filter_conditions.append(role_condition)
            
        if category:
            cat_condition = qmodels.FieldCondition(
                key="category",
                match=qmodels.MatchValue(value=category)
            )
            filter_conditions.append(cat_condition)
            
        query_filter = None
        if filter_conditions:
            query_filter = qmodels.Filter(must=filter_conditions)
            
        results = qdrant_client._client.search(
            collection_name=COLLECTION_NAME,
            query_vector=query_vector,
            query_filter=query_filter,
            limit=limit
        )
        formatted = []
        for r in results:
            payload = r.payload or {}
            formatted.append({
                "score": r.score,
                "content": payload.get("content", ""),
                "filename": payload.get("filename", ""),
                "category": payload.get("category", ""),
                "uploaded_by": payload.get("uploaded_by", ""),
                "allowed_roles": payload.get("allowed_roles", []),
                "version": payload.get("version", 1),
                "chunk_index": payload.get("chunk_index", 0)
            })
            
        return formatted

    @staticmethod
    async def delete_document(filename: str) -> bool:
        """Deletes all chunks of a document from the Qdrant index."""
        try:
            qdrant_client.delete(
                collection_name=COLLECTION_NAME,
                points_selector=qmodels.Filter(
                    must=[
                        qmodels.FieldCondition(
                            key="filename",
                            match=qmodels.MatchValue(value=filename)
                        )
                    ]
                )
            )
            return True
        except Exception:
            return False
